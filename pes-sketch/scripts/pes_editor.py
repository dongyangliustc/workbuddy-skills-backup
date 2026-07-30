#!/usr/bin/env python3
"""
PES Editor — 势能面 PPTX 全局参数批量调整工具 (Windows GUI)

功能：
  1. 批量修改 PPTX 中的所有连接线/横线宽度
  2. 批量修改 PPT 字体和字号
  3. 一键重新对齐连接线端点到最近的物种横线端点（模拟粘合效果）

用法：
  python pes_editor.py
  或直接双击运行。
"""

import sys
import os
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pptx import Presentation
from pptx.util import Pt, Emu, Inches
from pptx.oxml.ns import qn
import math


# ============================================================
# PPTX 处理引擎
# ============================================================

def scan_pptx(filepath):
    """扫描 PPTX，返回结构化信息"""
    prs = Presentation(filepath)
    info = {
        "slides": len(prs.slides),
        "width": prs.slide_width,
        "height": prs.slide_height,
        "items": [],  # (slide_idx, shape_type, subtype, current_value, shape)
    }

    for slide_idx, slide in enumerate(prs.slides):
        for el in _iter_slide_elements(slide):
            tag = el.tag.split('}')[-1]

            # ---- 连接器（横线、轴线、连接线） ----
            if tag == 'cxnSp':
                ln = el.find('.//' + qn('a:ln'))
                if ln is None:
                    continue
                w_emu = ln.get('w')
                w_pt = round(int(w_emu) / 12700, 1) if w_emu else 0
                is_dashed = el.find('.//' + qn('a:prstDash')) is not None

                # 判断是否为物种横线
                if not is_dashed:
                    xfrm = el.find('.//' + qn('a:xfrm'))
                    if xfrm is not None:
                        ext = xfrm.find(qn('a:ext'))
                        off = xfrm.find(qn('a:off'))
                        if ext is not None and off is not None:
                            cx = int(ext.get('cx', '0'))
                            cy = int(ext.get('cy', '0'))
                            x = int(off.get('x', '0'))
                            # 短水平线 → 物种横线（长度 > 1.2cm）
                            if cy < 5000 and cx > 400000:
                                info["items"].append(
                                    (slide_idx, "connector", "marker", w_pt, el)
                                )
                                continue
                    # 否则为轴线/刻度线
                    info["items"].append(
                        (slide_idx, "connector", "axis", w_pt, el)
                    )
                else:
                    # 虚线 → 连接线
                    info["items"].append(
                        (slide_idx, "connector", "dashed", w_pt, el)
                    )

            # ---- 文本框（同时检查 rPr 和 defRPr） ----
            elif tag == 'sp':
                txBody = el.find(qn('p:txBody'))
                if txBody is not None:
                    for rPr in txBody.iter(qn('a:rPr')):
                        _record_font_info(slide_idx, rPr, info, el)
                    for defRPr in txBody.iter(qn('a:defRPr')):
                        _record_font_info(slide_idx, defRPr, info, el)

    return info


def _record_font_info(slide_idx, rPr, info, el):
    """记录一个 rPr/defRPr 元素的字体信息"""
    sz = rPr.get('sz')
    sz_pt = round(int(sz) / 100, 1) if sz else 11
    font_family = None
    latin = rPr.find(qn('a:latin'))
    if latin is not None:
        font_family = latin.get('typeface')
    if font_family is None:
        ea = rPr.find(qn('a:ea'))
        if ea is not None:
            font_family = ea.get('typeface')
    info["items"].append(
        (slide_idx, "text", "text", (font_family or "?", sz_pt), el)
    )


def _iter_slide_elements(slide):
    """递归遍历幻灯片中所有形状元素（包括组内的）"""
    for shape in slide.shapes:
        el = shape._element
        tag = el.tag.split('}')[-1]
        if tag == 'grpSp':
            for child in el:
                ctag = child.tag.split('}')[-1]
                if ctag in ('sp', 'cxnSp', 'pic'):
                    yield child
        elif tag in ('sp', 'cxnSp', 'pic'):
            yield el


def set_marker_width(prs, target_pt):
    """修改所有物种横线宽度"""
    modified = 0
    target_emu = int(target_pt * 12700)
    for slide in prs.slides:
        for el in _iter_slide_elements(slide):
            if el.tag != qn('p:cxnSp'):
                continue
            is_dashed = el.find('.//' + qn('a:prstDash')) is not None
            if is_dashed:
                continue
            xfrm = el.find('.//' + qn('a:xfrm'))
            if xfrm is not None:
                ext = xfrm.find(qn('a:ext'))
                if ext is not None:
                    cy = int(ext.get('cy', '0'))
                    cx = int(ext.get('cx', '0'))
                    if cy < 5000 and cx > 400000:
                        ln = el.find('.//' + qn('a:ln'))
                        if ln is not None:
                            ln.set('w', str(target_emu))
                            modified += 1
    return modified


def set_dashed_width(prs, target_pt):
    """修改所有虚线连接线宽度"""
    modified = 0
    target_emu = int(target_pt * 12700)
    for slide in prs.slides:
        for el in _iter_slide_elements(slide):
            if el.tag != qn('p:cxnSp'):
                continue
            is_dashed = el.find('.//' + qn('a:prstDash')) is not None
            if not is_dashed:
                continue
            ln = el.find('.//' + qn('a:ln'))
            if ln is not None:
                ln.set('w', str(target_emu))
                modified += 1
    return modified


def set_font(prs, font_name=None, font_size_pt=None):
    """批量修改字体和字号"""
    modified = 0
    if font_size_pt is not None:
        target_sz = str(int(font_size_pt * 100))

    for slide in prs.slides:
        for el in _iter_slide_elements(slide):
            # 同时处理 rPr 和 defRPr
            for rPr in list(el.iter(qn('a:rPr'))) + list(el.iter(qn('a:defRPr'))):
                changed = False
                if font_name:
                    from lxml import etree
                    latin = rPr.find(qn('a:latin'))
                    if latin is None:
                        latin = etree.SubElement(rPr, qn('a:latin'))
                    latin.set('typeface', font_name)
                    ea = rPr.find(qn('a:ea'))
                    if ea is None:
                        ea = etree.SubElement(rPr, qn('a:ea'))
                    ea.set('typeface', font_name)
                    changed = True
                if font_size_pt is not None:
                    rPr.set('sz', target_sz)
                    changed = True
                if changed:
                    modified += 1

    return modified


def re_align_connectors(prs):
    """
    重新对齐连接线端点：检测物种横线（短水平实线）和虚线连接线，
    将每条虚线连接线的两端对齐到最近的横线端点。
    """
    for slide in prs.slides:
        # 第一步：收集本页所有物种横线的端点
        markers = []
        for el in _iter_slide_elements(slide):
            if el.tag != qn('p:cxnSp'):
                continue
            is_dashed = el.find('.//' + qn('a:prstDash')) is not None
            if is_dashed:
                continue
            xfrm = el.find('.//' + qn('a:xfrm'))
            if xfrm is None:
                continue
            off = xfrm.find(qn('a:off'))
            ext = xfrm.find(qn('a:ext'))
            if off is None or ext is None:
                continue
            cx = int(ext.get('cx', '0'))
            cy = int(ext.get('cy', '0'))
            if cy < 5000 and cx > 400000:
                x = int(off.get('x', '0'))
                y = int(off.get('y', '0'))
                markers.append((x, x + cx, y))

        if len(markers) < 2:
            continue

        # 第二步：调整每条虚线连接线的端点
        for el in _iter_slide_elements(slide):
            if el.tag != qn('p:cxnSp'):
                continue
            is_dashed = el.find('.//' + qn('a:prstDash')) is not None
            if not is_dashed:
                continue

            # 获取连接线当前端点（通过 xfrm + flip 推断）
            xfrm = el.find('.//' + qn('a:xfrm'))
            if xfrm is None:
                continue
            off = xfrm.find(qn('a:off'))
            ext = xfrm.find(qn('a:ext'))
            if off is None or ext is None:
                continue
            cx = int(ext.get('cx', '0'))
            cy = int(ext.get('cy', '0'))
            x0 = int(off.get('x', '0'))
            y0 = int(off.get('y', '0'))
            flipH = xfrm.get('flipH')
            flipV = xfrm.get('flipV')

            # 推断连接线的两个端点
            # 对于直线连接器，端点在 (x0,y0) 和 (x0+cx, y0+cy)
            # flipH/flipV 影响端点顺序
            x1, y1 = x0, y0
            x2, y2 = x0 + cx, y0 + cy

            # 找到最近的横线端点
            def find_nearest_marker_end(tx, ty):
                best = None
                best_dist = float('inf')
                for ml, mr, my in markers:
                    for mx in (ml, mr):
                        d = math.sqrt((mx - tx) ** 2 + (my - ty) ** 2)
                        if d < best_dist:
                            best_dist = d
                            best = (mx, my, (ml, mr, my))
                return best, best_dist

            (nx1, ny1, _), d1 = find_nearest_marker_end(x1, y1)
            (nx2, ny2, _), d2 = find_nearest_marker_end(x2, y2)

            # 只有当距离合理时才调整
            max_dist = 500000  # 约 1.3cm，防止误连到远处的横线
            if d1 > max_dist or d2 > max_dist:
                continue

            # 重新计算 xfrm
            new_x = min(nx1, nx2)
            new_y = min(ny1, ny2)
            new_cx = abs(nx2 - nx1)
            new_cy = abs(ny2 - ny1)

            off.set('x', str(new_x))
            off.set('y', str(new_y))
            ext.set('cx', str(new_cx))
            ext.set('cy', str(new_cy))

            # 更新翻转标志
            if (nx1 > nx2) != (flipH == '1'):
                if nx1 > nx2:
                    xfrm.set('flipH', '1')
                else:
                    try:
                        del xfrm.attrib[qn('flipH' if 'flipH' in str(xfrm.attrib) else '')]
                    except:
                        pass
            if (ny1 > ny2) != (flipV == '1'):
                if ny1 > ny2:
                    xfrm.set('flipV', '1')
                else:
                    try:
                        del xfrm.attrib[qn('flipV' if 'flipV' in str(xfrm.attrib) else '')]
                    except:
                        pass


def apply_all(prs, config):
    """综合应用所有参数修改"""
    results = {}
    if config.get("marker_width"):
        results["marker_width"] = set_marker_width(prs, config["marker_width"])
    if config.get("dashed_width"):
        results["dashed_width"] = set_dashed_width(prs, config["dashed_width"])
    if config.get("font_name") or config.get("font_size"):
        results["font"] = set_font(prs, config.get("font_name"), config.get("font_size"))
    return results


# ============================================================
# GUI 界面
# ============================================================

class PESEditorApp:
    def __init__(self, root):
        self.root = root
        self.root.title("PES 全局参数调整工具")
        self.root.geometry("580x520")
        self.root.resizable(False, False)

        self.pptx_path = tk.StringVar()
        self.font_var = tk.StringVar(value="Times New Roman")
        self.font_size_var = tk.StringVar(value="11")
        self.marker_width_var = tk.StringVar(value="3")
        self.dashed_width_var = tk.StringVar(value="3")

        self._build_ui()

    def _build_ui(self):
        main_frame = ttk.Frame(self.root, padding=16)
        main_frame.pack(fill=tk.BOTH, expand=True)

        # ---- 标题 ----
        title_label = ttk.Label(main_frame, text="PES 全局参数批量调整工具",
                                font=("Segoe UI", 14, "bold"))
        title_label.pack(anchor=tk.W, pady=(0, 12))

        # ---- 文件选择 ----
        file_frame = ttk.LabelFrame(main_frame, text="PPTX 文件", padding=8)
        file_frame.pack(fill=tk.X, pady=(0, 10))

        file_row = ttk.Frame(file_frame)
        file_row.pack(fill=tk.X)
        ttk.Entry(file_row, textvariable=self.pptx_path, width=50).pack(side=tk.LEFT, padx=(0, 6))
        ttk.Button(file_row, text="浏览...", command=self._browse_file, width=10).pack(side=tk.LEFT)

        # ---- 全局参数 ----
        param_frame = ttk.LabelFrame(main_frame, text="全局参数设置", padding=12)
        param_frame.pack(fill=tk.X, pady=(0, 10))

        grid = ttk.Frame(param_frame)
        grid.pack(fill=tk.X)

        # 字体
        ttk.Label(grid, text="字体:").grid(row=0, column=0, sticky=tk.W, padx=(0, 8), pady=4)
        font_combo = ttk.Combobox(grid, textvariable=self.font_var, width=20,
                                  values=["Times New Roman", "Arial", "Calibri",
                                          "宋体", "SimSun", "微软雅黑", "KaiTi",
                                          "Helvetica", "Courier New", "Symbol"])
        font_combo.grid(row=0, column=1, sticky=tk.W, pady=4)
        font_combo.bind('<<ComboboxSelected>>', lambda e: None)

        ttk.Label(grid, text="字号 (pt):").grid(row=0, column=2, sticky=tk.W, padx=(16, 8), pady=4)
        ttk.Entry(grid, textvariable=self.font_size_var, width=8).grid(row=0, column=3, sticky=tk.W, pady=4)

        # 横线宽度
        ttk.Label(grid, text="物种横线宽度 (pt):").grid(row=1, column=0, sticky=tk.W, padx=(0, 8), pady=4)
        ttk.Entry(grid, textvariable=self.marker_width_var, width=8).grid(row=1, column=1, sticky=tk.W, pady=4)

        ttk.Label(grid, text="连接线宽度 (pt):").grid(row=1, column=2, sticky=tk.W, padx=(16, 8), pady=4)
        ttk.Entry(grid, textvariable=self.dashed_width_var, width=8).grid(row=1, column=3, sticky=tk.W, pady=4)

        # ---- 按钮区 ----
        btn_frame = ttk.Frame(main_frame)
        btn_frame.pack(fill=tk.X, pady=(0, 10))

        ttk.Button(btn_frame, text="应用全局参数",
                   command=self._apply_params, width=18).pack(side=tk.LEFT, padx=(0, 8))

        ttk.Button(btn_frame, text="重新对齐连接线端点",
                   command=self._re_align, width=20).pack(side=tk.LEFT)

        # ---- 状态栏 ----
        status_frame = ttk.LabelFrame(main_frame, text="状态", padding=6)
        status_frame.pack(fill=tk.BOTH, expand=True)

        self.status_text = tk.Text(status_frame, height=6, wrap=tk.WORD,
                                   font=("Consolas", 9), bg="#fafafa", relief=tk.FLAT)
        self.status_text.pack(fill=tk.BOTH, expand=True)

        scrollbar = ttk.Scrollbar(self.status_text, command=self.status_text.yview)
        scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        self.status_text.config(yscrollcommand=scrollbar.set)

        self._log("就绪。请选择 PPTX 文件后操作。")

    def _log(self, msg):
        self.status_text.insert(tk.END, msg + "\n")
        self.status_text.see(tk.END)
        self.root.update_idletasks()

    def _browse_file(self):
        path = filedialog.askopenfilename(
            title="选择 PPTX 文件",
            filetypes=[("PowerPoint 文件", "*.pptx"), ("All Files", "*.*")]
        )
        if path:
            self.pptx_path.set(path)
            self._log(f"已选择: {os.path.basename(path)}")
            try:
                info = scan_pptx(path)
                self._log(f"  幻灯片数: {info['slides']}")
                markers = sum(1 for i in info["items"] if i[2] == "marker")
                dashed = sum(1 for i in info["items"] if i[2] == "dashed")
                texts = sum(1 for i in info["items"] if i[1] == "text")
                self._log(f"  物种横线: {markers} 条 | 虚线连接: {dashed} 条 | 文本框: {texts} 个")
            except Exception as e:
                self._log(f"  [ERR] 读取失败: {e}")

    def _get_ppt_path(self):
        path = self.pptx_path.get().strip()
        if not path:
            messagebox.showwarning("提示", "请先选择 PPTX 文件")
            return None
        if not os.path.isfile(path):
            messagebox.showerror("错误", f"文件不存在: {path}")
            return None
        return path

    def _apply_params(self):
        path = self._get_ppt_path()
        if not path:
            return

        config = {}
        try:
            if self.font_var.get().strip():
                config["font_name"] = self.font_var.get().strip()
            if self.font_size_var.get().strip():
                config["font_size"] = float(self.font_size_var.get().strip())
            if self.marker_width_var.get().strip():
                config["marker_width"] = float(self.marker_width_var.get().strip())
            if self.dashed_width_var.get().strip():
                config["dashed_width"] = float(self.dashed_width_var.get().strip())
        except ValueError:
            messagebox.showerror("错误", "请输入有效的数值（字号、宽度）")
            return

        if not config:
            messagebox.showwarning("提示", "至少填写一项参数")
            return

        try:
            prs = Presentation(path)
            results = apply_all(prs, config)

            # 另存为新文件
            base, ext = os.path.splitext(path)
            out_path = f"{base}_modified{ext}"
            prs.save(out_path)

            self._log(f"--- 应用全局参数 ---")
            if "marker_width" in results:
                self._log(f"  物种横线: {results['marker_width']} 条 → {config['marker_width']}pt")
            if "dashed_width" in results:
                self._log(f"  虚线连接: {results['dashed_width']} 条 → {config['dashed_width']}pt")
            if "font" in results:
                font_msg = ""
                if "font_name" in config:
                    font_msg += f"字体={config['font_name']} "
                if "font_size" in config:
                    font_msg += f"字号={config['font_size']}pt"
                self._log(f"  文本框: {results['font']} 处设置 → {font_msg}")

            self._log(f"  已保存: {os.path.basename(out_path)}")
            self._log(f"  (原始文件未修改)")

        except Exception as e:
            self._log(f"  [ERR] 操作失败: {e}")
            import traceback
            self._log(traceback.format_exc())

    def _re_align(self):
        path = self._get_ppt_path()
        if not path:
            return

        try:
            prs = Presentation(path)
            re_align_connectors(prs)

            base, ext = os.path.splitext(path)
            out_path = f"{base}_aligned{ext}"
            prs.save(out_path)

            self._log(f"--- 重新对齐连接线端点 ---")
            self._log(f"  已将虚线连接线端点对齐到最近的物种横线端点")
            self._log(f"  已保存: {os.path.basename(out_path)}")
            self._log(f"  (提示：在 PPT 中取消组合后移动横线，连接线将保持对齐)")

        except Exception as e:
            self._log(f"  [ERR] 对齐失败: {e}")
            import traceback
            self._log(traceback.format_exc())


# ============================================================
# 入口
# ============================================================

if __name__ == "__main__":
    root = tk.Tk()
    app = PESEditorApp(root)
    root.mainloop()
