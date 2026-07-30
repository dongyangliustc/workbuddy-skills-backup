#!/usr/bin/env python3
"""
PES Resizer — 势能面 PPTX 线宽批量调整工具

一键修改已生成的 PES PPTX 文件中所有连接线的宽度。
支持按模式筛选（物种横线 / 虚线连接线 / 全部线）。

Usage:
    # 查看当前 PPTX 中所有连接线的宽度
    python pes_resizer.py input.pptx --status

    # 将所有连接线宽度设为 3pt
    python pes_resizer.py input.pptx --set 3

    # 仅修改虚线（连接线），设为 2pt
    python pes_resizer.py input.pptx --set 2 --mode dashed

    # 仅修改实线（物种横线、轴线），设为 4pt
    python pes_resizer.py input.pptx --set 4 --mode solid

    # 仅修改指定颜色的线（按十六进制 RGB）
    python pes_resizer.py input.pptx --set 3 --color 333333

    # 全部设为 2pt 并另存为新文件
    python pes_resizer.py input.pptx --set 2 -o modified.pptx
"""

import sys
import os
import argparse
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn


def emu_to_pt(emu_val):
    """EMU 转换为 pt"""
    if emu_val is None:
        return None
    return int(emu_val) / 12700.0


def pt_to_emu(pt_val):
    """pt 转换为 EMU"""
    return int(pt_val * 12700)


def format_color(shape):
    """获取连接线的颜色（十六进制字符串）"""
    try:
        ln = shape._element.find(qn('a:ln'))
        if ln is not None:
            solidFill = ln.find(qn('a:solidFill'))
            if solidFill is not None:
                srgb = solidFill.find(qn('a:srgbClr'))
                if srgb is not None:
                    return srgb.get('val')
    except Exception:
        pass
    return None


def is_dashed(shape):
    """判断连接线是否为虚线"""
    try:
        # 先尝试在 spPr 下找 ln
        spPr = shape._element.find(qn('p:spPr'))
        if spPr is not None:
            ln = spPr.find(qn('a:ln'))
            if ln is not None:
                prstDash = ln.find(qn('a:prstDash'))
                return prstDash is not None
        # 再尝试直接在 cxnSp 下找 ln
        ln = shape._element.find(qn('a:ln'))
        if ln is not None:
            prstDash = ln.find(qn('a:prstDash'))
            return prstDash is not None
    except Exception:
        pass
    return False


def is_species_marker(shape, slide_width):
    """启发式判断是否为物种横线标记（短横线）"""
    # 物种横线是连接器，且起点和终点的 y 坐标相同（水平）
    # 且长度通常接近 SPECIES_MARKER_LENGTH
    try:
        cxnSp = shape._element
        spPr = cxnSp.find(qn('p:cxnSpPr')) if cxnSp.find(qn('p:cxnSpPr')) is not None else cxnSp.find(qn('a:xfrm'))
        # 获取连接器端点
        for child in cxnSp.iter():
            if child.tag == qn('a:cxn'):
                x1 = child.get('x')
                # Check if it's a connector with flips
                break
        # 检查是否水平
        xfrm = cxnSp.find(qn('p:cxnSpPr'))
        if xfrm is None:
            xfrm = cxnSp.find(qn('a:xfrm'))
        if xfrm is not None:
            off = xfrm.find(qn('a:off'))
            ext = xfrm.find(qn('a:ext'))
            if off is not None and ext is not None:
                h = int(ext.get('y', '0'))
                if h < 10000:  # 近乎水平
                    return True
    except Exception:
        pass
    return False


def scan_connectors(prs):
    """扫描 PPTX 中所有连接线，返回统计信息"""
    connectors = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            # 只有连接器（MSO_SHAPE_TYPE.CONNECTOR = 9）
            # cxnSp 在 XML 层面判断更可靠
            if shape._element.tag == qn('p:cxnSp'):
                try:
                    ln = shape._element.find(qn('a:ln'))
                    if ln is None:
                        # 有时 ln 在 spPr 下
                        spPr = shape._element.find(qn('p:spPr'))
                        if spPr is not None:
                            ln = spPr.find(qn('a:ln'))
                except:
                    ln = None

                w_emu = None
                if ln is not None:
                    w_emu = ln.get('w')
                
                w_pt = emu_to_pt(w_emu) if w_emu else None
                dashed = is_dashed(shape)
                color = format_color(shape)
                
                connectors.append({
                    'slide': slide_idx + 1,
                    'name': shape.name or '',
                    'width_emu': w_emu,
                    'width_pt': round(w_pt, 1) if w_pt else None,
                    'dashed': dashed,
                    'color': color,
                })
    return connectors


def set_line_width(prs, target_pt, mode='all', target_color=None):
    """批量修改连接线宽度
    
    Args:
        prs: Presentation 对象
        target_pt: 目标线宽（pt）
        mode: 'all' | 'dashed' | 'solid'
        target_color: 可选，仅修改指定颜色的线（如 '333333'）
    """
    target_emu = pt_to_emu(target_pt)
    modified = 0
    skipped = 0

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape._element.tag != qn('p:cxnSp'):
                continue

            try:
                spPr = shape._element.find(qn('p:spPr'))
                if spPr is None:
                    spPr = shape._element
                ln = spPr.find(qn('a:ln'))
                # 如果没有 ln 元素，创建一个
                if ln is None:
                    # 可能需要创建 ln 元素
                    continue
                
                # 模式过滤
                dashed = False
                prstDash = ln.find(qn('a:prstDash'))
                if prstDash is not None:
                    dashed = True
                
                if mode == 'dashed' and not dashed:
                    skipped += 1
                    continue
                if mode == 'solid' and dashed:
                    skipped += 1
                    continue

                # 颜色过滤
                if target_color is not None:
                    solidFill = ln.find(qn('a:solidFill'))
                    if solidFill is not None:
                        srgb = solidFill.find(qn('a:srgbClr'))
                        color_val = srgb.get('val') if srgb is not None else None
                        if color_val != target_color:
                            skipped += 1
                            continue

                # 修改宽度
                ln.set('w', str(target_emu))
                modified += 1

            except Exception as e:
                print(f"  [WARN] 跳过形状 '{getattr(shape, 'name', '?')}': {e}")
                skipped += 1

    return modified, skipped


def main():
    parser = argparse.ArgumentParser(
        description="PES Resizer — 批量调整势能面 PPTX 连接线宽度"
    )
    parser.add_argument("input", help="输入 PPTX 文件路径")
    parser.add_argument("-o", "--output", help="输出 PPTX 文件路径（默认覆盖原文件）")
    parser.add_argument("--set", type=float, help="目标线宽（pt）")
    parser.add_argument("--mode", choices=["all", "dashed", "solid"], default="all",
                        help="筛选模式: all=全部, dashed=仅虚线, solid=仅实线")
    parser.add_argument("--color", help="按颜色筛选（十六进制，如 333333）")
    parser.add_argument("--status", action="store_true", help="仅查看当前宽度统计，不做修改")
    args = parser.parse_args()

    if not os.path.isfile(args.input):
        print(f"[ERR] 文件不存在: {args.input}")
        sys.exit(1)

    if not args.status and args.set is None:
        print("[ERR] 请指定 --set <pt> 或 --status 查看当前状态")
        sys.exit(1)

    prs = Presentation(args.input)

    # 扫描
    connectors = scan_connectors(prs)
    
    if not connectors:
        print("[INFO] 未找到连接线")
        return

    print(f"=== 当前连接线统计 ===")
    print(f"幻灯片数: {len(prs.slides)}")
    print(f"连接线总数: {len(connectors)}")

    # 统计分组
    dashed_count = sum(1 for c in connectors if c['dashed'])
    solid_count = sum(1 for c in connectors if not c['dashed'])
    print(f"  实线: {solid_count}")
    print(f"  虚线: {dashed_count}")

    widths = [c['width_pt'] for c in connectors if c['width_pt'] is not None]
    if widths:
        print(f"  宽度范围: {min(widths):.1f}pt ~ {max(widths):.1f}pt")
        from collections import Counter
        wc = Counter(widths)
        print(f"  宽度分布:")
        for w, cnt in sorted(wc.items()):
            dashed_in_group = sum(1 for c in connectors if c['width_pt'] == w and c['dashed'])
            solid_in_group = sum(1 for c in connectors if c['width_pt'] == w and not c['dashed'])
            label = "dashed" if dashed_in_group > solid_in_group else "solid"
            print(f"    {w:.1f}pt → {cnt} 条 ({'虚线' if dashed_in_group > solid_in_group else '实线'})")

    if args.status:
        return

    # 执行修改
    target_pt = args.set
    target_color = args.color
    print(f"\n=== 执行修改 ===")
    print(f"目标宽度: {target_pt:.1f}pt")
    print(f"筛选模式: {args.mode}")
    if target_color:
        print(f"颜色筛选: #{target_color}")

    modified, skipped = set_line_width(prs, target_pt, args.mode, target_color)
    print(f"已修改: {modified} 条")
    print(f"已跳过: {skipped} 条")

    if modified == 0:
        print("[WARN] 没有连接线被修改")
        return

    # 保存
    output_path = args.output or args.input
    prs.save(output_path)
    print(f"[OK] 已保存: {output_path}")


if __name__ == "__main__":
    main()
