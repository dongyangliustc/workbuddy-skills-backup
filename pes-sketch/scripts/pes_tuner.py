#!/usr/bin/env python3
"""
PES Tuner — 势能面 PPTX 全局参数直观调整工具 (Windows 桌面版)

用法：双击运行，或 python pes_tuner.py
"""

import sys
import os
import tkinter as tk
from tkinter import ttk, filedialog, messagebox
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn
import math


# ============================================================
# 核心引擎
# ============================================================

def _iter_slide_elements(slide):
    """递归遍历幻灯片中所有形状元素（包括组内的）"""
    for shape in slide.shapes:
        el = shape._element
        tag = el.tag.split('}')[-1]
        if tag == 'grpSp':
            for child in el:
                ctag = child.tag.split('}')[-1]
                if ctag in ('sp', 'cxnSp', 'pic'):
                    yield child
        elif tag in ('sp', 'cxnSp', 'pic'):
            yield el


def scan_pptx(filepath):
    """扫描 PPTX，返回统计信息"""
    prs = Presentation(filepath)
    info = {"slides": len(prs.slides), "markers": 0, "dashed": 0, "texts": 0}

    for slide in prs.slides:
        for el in _iter_slide_elements(slide):
            tag = el.tag.split('}')[-1]
            if tag == 'cxnSp':
                is_dashed = el.find('.//' + qn('a:prstDash')) is not None
                if is_dashed:
                    info["dashed"] += 1
                else:
                    xfrm = el.find('.//' + qn('a:xfrm'))
                    if xfrm is not None:
                        ext = xfrm.find(qn('a:ext'))
                        if ext is not None:
                            cy = int(ext.get('cy', '0'))
                            cx = int(ext.get('cx', '0'))
                            if cy < 5000 and cx > 400000:
                                info["markers"] += 1
            elif tag == 'sp':
                txBody = el.find(qn('p:txBody'))
                if txBody is not None:
                    for _ in txBody.iter(qn('a:defRPr')):
                        info["texts"] += 1
                    for _ in txBody.iter(qn('a:rPr')):
                        info["texts"] += 1

    return info

def set_marker_width(prs, target_pt):
    cnt = 0
    target_emu = int(target_pt * 12700)
    for slide in prs.slides:
        for el in _iter_slide_elements(slide):
            if el.tag != qn('p:cxnSp'):
                continue
            if el.find('.//' + qn('a:prstDash')) is not None:
                continue
            xfrm = el.find('.//' + qn('a:xfrm'))
            if xfrm is None:
                continue
            ext = xfrm.find(qn('a:ext'))
            if ext is None:
                continue
            cy = int(ext.get('cy', '0'))
            cx = int(ext.get('cx', '0'))
            if cy < 5000 and cx > 400000:
                ln = el.find('.//' + qn('a:ln'))
                if ln is not None:
                    ln.set('w', str(target_emu))
                    cnt += 1
    return cnt

def set_dashed_width(prs, target_pt):
    cnt = 0
    target_emu = int(target_pt * 12700)
    for slide in prs.slides:
        for el in _iter_slide_elements(slide):
            if el.tag != qn('p:cxnSp'):
                continue
            if el.find('.//' + qn('a:prstDash')) is None:
                continue
            ln = el.find('.//' + qn('a:ln'))
            if ln is not None:
                ln.set('w', str(target_emu))
                cnt += 1
    return cnt

def set_font(prs, font_name=None, font_size_pt=None):
    cnt = 0
    target_sz = str(int(font_size_pt * 100)) if font_size_pt else None
    for slide in prs.slides:
        for el in _iter_slide_elements(slide):
            for rPr in list(el.iter(qn('a:rPr'))) + list(el.iter(qn('a:defRPr'))):
                changed = False
                if font_name:
                    from lxml import etree
                    latin = rPr.find(qn('a:latin'))
                    if latin is None:
                        latin = etree.SubElement(rPr, qn('a:latin'))
                    latin.set('typeface', font_name)
                    ea = rPr.find(qn('a:ea'))
                    if ea is None:
                        ea = etree.SubElement(rPr, qn('a:ea'))
                    ea.set('typeface', font_name)
                    changed = True
                if target_sz:
                    rPr.set('sz', target_sz)
                    changed = True
                if changed:
                    cnt += 1
    return cnt

def re_align_connectors(prs):
    for slide in prs.slides:
        markers = []
        for el in _iter_slide_elements(slide):
            if el.tag != qn('p:cxnSp'):
                continue
            if el.find('.//' + qn('a:prstDash')) is not None:
                continue
            xfrm = el.find('.//' + qn('a:xfrm'))
            if xfrm is None:
                continue
            off = xfrm.find(qn('a:off'))
            ext = xfrm.find(qn('a:ext'))
            if off is None or ext is None:
                continue
            cx = int(ext.get('cx', '0'))
            cy = int(ext.get('cy', '0'))
            if cy < 5000 and cx > 400000:
                x = int(off.get('x', '0'))
                y = int(off.get('y', '0'))
                markers.append((x, x + cx, y))

        if len(markers) < 2:
            continue

        for el in _iter_slide_elements(slide):
            if el.tag != qn('p:cxnSp'):
                continue
            if el.find('.//' + qn('a:prstDash')) is None:
                continue
            xfrm = el.find('.//' + qn('a:xfrm'))
            if xfrm is None:
                continue
            off = xfrm.find(qn('a:off'))
            ext = xfrm.find(qn('a:ext'))
            if off is None or ext is None:
                continue
            cx = int(ext.get('cx', '0'))
            cy = int(ext.get('cy', '0'))
            x0 = int(off.get('x', '0'))
            y0 = int(off.get('y', '0'))
            x1, y1 = x0, y0
            x2, y2 = x0 + cx, y0 + cy

            def _nearest(tx, ty):
                best, bd = None, float('inf')
                for ml, mr, my in markers:
                    for mx in (ml, mr):
                        d = math.sqrt((mx - tx) ** 2 + (my - ty) ** 2)
                        if d < bd:
                            bd = d
                            best = (mx, my)
                return best, bd

            (nx1, ny1), d1 = _nearest(x1, y1)
            (nx2, ny2), d2 = _nearest(x2, y2)
            if d1 > 500000 or d2 > 500000:
                continue

            off.set('x', str(min(nx1, nx2)))
            off.set('y', str(min(ny1, ny2)))
            ext.set('cx', str(abs(nx2 - nx1)))
            ext.set('cy', str(abs(ny2 - ny1)))


# ============================================================
# GUI (两列布局：参数 | 值)
# ============================================================

class PESTuner:
    def __init__(self, root):
        self.root = root
        self.root.title("PES Tuner — 势能面参数调整")
        self.root.geometry("520x420")
        self.root.resizable(False, False)
        self.pptx_path = None
        self.prs = None

        self._build_ui()

    def _build_ui(self):
        # ---- 顶部：文件选择 ----
        top = ttk.Frame(self.root, padding=8)
        top.pack(fill=tk.X)
        ttk.Label(top, text="PPTX 文件：").pack(side=tk.LEFT)
        self.file_label = ttk.Label(top, text="（未选择）", foreground="#888",
                                     width=30, anchor=tk.W)
        self.file_label.pack(side=tk.LEFT, padx=4)
        ttk.Button(top, text="浏览", command=self._browse,
                   width=8).pack(side=tk.LEFT, padx=4)
        self.scan_label = ttk.Label(top, text="", foreground="#666",
                                     font=("", 8))
        self.scan_label.pack(side=tk.RIGHT)

        # ---- 主体：两列参数表 ----
        main = ttk.Frame(self.root, padding=(12, 4))
        main.pack(fill=tk.BOTH, expand=True)

        # 表头
        hdr = ttk.Frame(main)
        hdr.pack(fill=tk.X, pady=(0, 4))
        ttk.Label(hdr, text="参数", font=("", 10, "bold"),
                  width=22).pack(side=tk.LEFT)
        ttk.Label(hdr, text="当前值", font=("", 10, "bold"),
                  width=12).pack(side=tk.LEFT)
        ttk.Label(hdr, text="", font=("", 10, "bold"),
                  width=14).pack(side=tk.LEFT)

        # 分隔线
        ttk.Separator(main, orient=tk.HORIZONTAL).pack(fill=tk.X, pady=2)

        # 参数行容器
        rows = ttk.Frame(main)
        rows.pack(fill=tk.BOTH, expand=True, pady=4)

        def _make_row(label, widget_type, default="", choices=None, width=12):
            f = ttk.Frame(rows)
            f.pack(fill=tk.X, pady=3)
            ttk.Label(f, text=label, width=22, anchor=tk.W).pack(side=tk.LEFT)
            if widget_type == "combo":
                var = tk.StringVar(value=default)
                w = ttk.Combobox(f, textvariable=var, values=choices,
                                 width=width-2, state="readonly")
                w.pack(side=tk.LEFT)
            else:
                var = tk.StringVar(value=default)
                w = ttk.Entry(f, textvariable=var, width=width)
                w.pack(side=tk.LEFT)
            ttk.Label(f, text="", width=14).pack(side=tk.LEFT)
            return var

        self.var_font = _make_row("字体", "combo", "Times New Roman",
                                  ["Times New Roman", "Arial", "Calibri",
                                   "宋体", "SimSun", "微软雅黑", "KaiTi",
                                   "Courier New", "Symbol"])
        self.var_size = _make_row("字号 (pt)", "entry", "11")
        self.var_mw = _make_row("物种横线宽度 (pt)", "entry", "3")
        self.var_dw = _make_row("连接线宽度 (pt)", "entry", "3")

        # 按钮区
        action_frame = ttk.Frame(main)
        action_frame.pack(fill=tk.X, pady=(8, 2))
        
        self.btn_apply = ttk.Button(action_frame, text="应用参数",
                                     command=self._apply, width=18)
        self.btn_apply.pack(side=tk.LEFT, padx=(0, 8))
        self.btn_align = ttk.Button(action_frame, text="重新对齐连接线",
                                     command=self._align, width=18)
        self.btn_align.pack(side=tk.LEFT)
        
        # 覆盖原文件选项
        self.overwrite_var = tk.BooleanVar(value=True)
        ttk.Checkbutton(action_frame, text="覆盖原文件",
                        variable=self.overwrite_var).pack(side=tk.RIGHT)
        
        # 初始禁用
        self.btn_apply.state(["disabled"])
        self.btn_align.state(["disabled"])

        # ---- 底部状态 ----
        status_frame = ttk.Frame(self.root, padding=(8, 2))
        status_frame.pack(fill=tk.X)
        self.status = ttk.Label(status_frame, text="就绪 — 勾选「覆盖原文件」可直接保存到原文件",
                                 foreground="#555")
        self.status.pack(side=tk.LEFT)

    # ---- 业务方法 ----

    def _log(self, msg):
        self.status.config(text=msg)
        self.root.update_idletasks()

    def _browse(self):
        path = filedialog.askopenfilename(
            title="选择 PPTX 文件",
            filetypes=[("PowerPoint 文件", "*.pptx"), ("All Files", "*.*")]
        )
        if not path:
            return
        self.pptx_path = path
        self.file_label.config(text=os.path.basename(path), foreground="#000")
        try:
            info = scan_pptx(path)
            self.scan_label.config(
                text=f"横线:{info['markers']} 虚线:{info['dashed']} 文本:{info['texts']}")
            self._log(f"已加载：{os.path.basename(path)}  ({info['slides']} 页)")
            self.btn_apply.state(["!disabled"])
            self.btn_align.state(["!disabled"])
        except Exception as e:
            self._log(f"[错误] 读取失败：{e}")
            self.btn_apply.state(["disabled"])
            self.btn_align.state(["disabled"])

    def _apply(self):
        if not self.pptx_path:
            return
        try:
            font_name = self.var_font.get().strip() or None
            font_size = float(self.var_size.get().strip()) if self.var_size.get().strip() else None
            mw = float(self.var_mw.get().strip()) if self.var_mw.get().strip() else None
            dw = float(self.var_dw.get().strip()) if self.var_dw.get().strip() else None
        except ValueError:
            messagebox.showerror("错误", "字号和宽度请输入有效数字")
            return

        if not any([font_name, font_size, mw, dw]):
            messagebox.showwarning("提示", "至少填写一项参数")
            return

        try:
            prs = Presentation(self.pptx_path)
            parts = []
            if mw:
                c = set_marker_width(prs, mw)
                parts.append(f"横线={c}")
            if dw:
                c = set_dashed_width(prs, dw)
                parts.append(f"虚线={c}")
            if font_name or font_size:
                c = set_font(prs, font_name, font_size)
                parts.append(f"文本={c}")

            # 确定保存路径
            base, ext = os.path.splitext(self.pptx_path)
            if self.overwrite_var.get():
                out = self.pptx_path
            else:
                out = f"{base}_调整后{ext}"

            # 尝试保存（若文件被 PPT 锁定会抛出 PermissionError）
            try:
                prs.save(out)
                if self.overwrite_var.get():
                    self._log(f"✓ 已覆盖原文件 ({', '.join(parts)})")
                else:
                    self._log(f"✓ 已保存：{os.path.basename(out)}  ({', '.join(parts)})")
            except PermissionError:
                if self.overwrite_var.get():
                    # 覆盖失败，回退到 _调整后 后缀
                    fallback = f"{base}_调整后{ext}"
                    prs.save(fallback)
                    self._log(f"⚠ 原文件被 PPT 占用，无法覆盖")
                    self._log(f"  已保存为：{os.path.basename(fallback)}")
                    messagebox.showwarning(
                        "文件被占用",
                        f"「{os.path.basename(self.pptx_path)}」正在 PowerPoint 中打开，\n"
                        "无法直接覆盖。\n\n"
                        "已自动保存为「_调整后.pptx」版本。\n"
                        "请在 PPT 中关闭该文件后再试。"
                    )
                else:
                    raise
        except Exception as e:
            self._log(f"[错误] {e}")
            import traceback
            self._log(traceback.format_exc())

    def _align(self):
        if not self.pptx_path:
            return
        try:
            prs = Presentation(self.pptx_path)
            re_align_connectors(prs)

            base, ext = os.path.splitext(self.pptx_path)
            if self.overwrite_var.get():
                out = self.pptx_path
            else:
                out = f"{base}_对齐后{ext}"

            try:
                prs.save(out)
                self._log(f"✓ 连接线已{'覆盖' if self.overwrite_var.get() else '保存'} → {os.path.basename(out)}")
            except PermissionError:
                if self.overwrite_var.get():
                    fallback = f"{base}_对齐后{ext}"
                    prs.save(fallback)
                    self._log(f"⚠ 原文件被 PPT 占用，已保存为：{os.path.basename(fallback)}")
                    messagebox.showwarning(
                        "文件被占用",
                        f"「{os.path.basename(self.pptx_path)}」正在 PowerPoint 中打开，\n"
                        "无法直接覆盖。\n\n"
                        "已自动保存为「_对齐后.pptx」版本。"
                    )
                else:
                    raise
        except Exception as e:
            self._log(f"[错误] {e}")


# ============================================================
# 入口
# ============================================================

if __name__ == "__main__":
    root = tk.Tk()
    app = PESTuner(root)
    root.mainloop()
