#!/usr/bin/env python3
"""
PES Sketch — 反应势能面草图生成器 (PowerPoint)

从 YAML 配置文件生成可编辑的势能面 PPTX 文件。
支持单页完整视图和分步视图两种模式。

Usage:
    python pes_sketch.py config.yaml -o output.pptx
    python pes_sketch.py config.yaml -o output.pptx --mode single
    python pes_sketch.py config.yaml -o output.pptx --mode stepwise
"""

import sys
import os
import argparse
import yaml
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.dml import MSO_LINE_DASH_STYLE
import math


# ============================================================
# 颜色工具
# ============================================================
DEFAULT_WELL_COLOR = "#000000"    # 黑色
DEFAULT_BARRIER_COLOR = "#000000" # 黑色
DEFAULT_BIMOL_COLOR = "#000000"   # 黑色
AXIS_COLOR = "#333333"
LABEL_COLOR = "#222222"

# 默认线宽常量 (2 cm 转为 EMU)
# species 横线标记长度
SPECIES_MARKER_LENGTH = Emu(int(2 / 2.54 * 914400))  # 2cm 对应的 EMU
# 连接线宽度
CONNECTOR_LINE_WIDTH = Pt(3)  # 3pt 线宽，与 2cm 标记匹配
# 物种横线自身的线宽
SPECIES_MARKER_LINE_WIDTH = Pt(3)


def _remove_shadow(shape):
    """移除形状的阴影效果"""
    try:
        shape.shadow.inherit = False
    except Exception:
        pass


def _group_species_shapes(slide, group_map):
    """
    将每个物种的多个形状组合为一个 PowerPoint 组，方便用户后续调整。
    group_map: {species_name: [(shape, is_horizontal_marker_bool), ...]}
    其中 is_horizontal_marker_bool 标记是否为横线连接器（需要调整 cxn 坐标）
    """
    from lxml import etree
    from pptx.oxml.ns import qn
    import copy

    spTree = slide.shapes._spTree

    for sp_name, items in group_map.items():
        if len(items) < 2:
            continue

        shapes = [it[0] for it in items]
        is_markers = [it[1] for it in items]

        # 计算所有形状的包围盒
        min_x, min_y, max_x, max_y = None, None, None, None

        for shape in shapes:
            el = shape._element
            for xfrm in el.iter(qn('a:xfrm')):
                off = xfrm.find(qn('a:off'))
                ext = xfrm.find(qn('a:ext'))
                if off is not None and ext is not None:
                    x = int(off.get('x', '0'))
                    y = int(off.get('y', '0'))
                    cx = int(ext.get('cx', '0'))
                    cy = int(ext.get('cy', '0'))
                    if min_x is None or x < min_x:
                        min_x = x
                    if min_y is None or y < min_y:
                        min_y = y
                    if max_x is None or x + cx > max_x:
                        max_x = x + cx
                    if max_y is None or y + cy > max_y:
                        max_y = y + cy

            # 对于连接器，还需考虑 cxn 连接点的位置
            for cxn in el.iter(qn('a:cxn')):
                x = int(cxn.get('x', '0'))
                y = int(cxn.get('y', '0'))
                # 相对于 a:xfrm/a:off，需要加上 xfrm 偏移才是绝对坐标
                parent_xfrm = el.find('.//' + qn('a:xfrm'))
                if parent_xfrm is not None:
                    parent_off = parent_xfrm.find(qn('a:off'))
                    if parent_off is not None:
                        x += int(parent_off.get('x', '0'))
                        y += int(parent_off.get('y', '0'))
                if min_x is None or x < min_x:
                    min_x = x
                if max_x is None or x > max_x:
                    max_x = x
                if min_y is None or y < min_y:
                    min_y = y
                if max_y is None or y > max_y:
                    max_y = y

        if min_x is None:
            continue

        grp_w = max_x - min_x
        grp_h = max_y - min_y

        # 创建 <p:grpSp> 元素
        grpSp_ns = 'http://schemas.openxmlformats.org/presentationml/2006/main'
        a_ns = 'http://schemas.openxmlformats.org/drawingml/2006/main'
        NSMAP = {'p': grpSp_ns, 'a': a_ns, 'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'}

        grpSp = etree.Element(qn('p:grpSp'), nsmap=NSMAP)

        # nvGrpSpPr
        nv = etree.SubElement(grpSp, qn('p:nvGrpSpPr'))
        cnvPr = etree.SubElement(nv, qn('p:cNvPr'))
        cnvPr.set('id', '0')
        cnvPr.set('name', f'Species {sp_name}')
        nv.append(etree.SubElement(nv, qn('p:cNvGrpSpPr')))
        nv.append(etree.SubElement(nv, qn('p:nvPr')))

        # grpSpPr
        gsp = etree.SubElement(grpSp, qn('p:grpSpPr'))
        gxfrm = etree.SubElement(gsp, qn('a:xfrm'))
        goff = etree.SubElement(gxfrm, qn('a:off'))
        goff.set('x', str(min_x))
        goff.set('y', str(min_y))
        gext = etree.SubElement(gxfrm, qn('a:ext'))
        gext.set('cx', str(grp_w))
        gext.set('cy', str(grp_h))
        chOff = etree.SubElement(gxfrm, qn('a:chOff'))
        chOff.set('x', '0')
        chOff.set('y', '0')
        chExt = etree.SubElement(gxfrm, qn('a:chExt'))
        chExt.set('cx', str(grp_w))
        chExt.set('cy', str(grp_h))

        # 移动形状到组中，调整坐标
        for shape in shapes:
            el = shape._element
            # 调整 a:xfrm/a:off，减去组原点
            for xfrm in el.iter(qn('a:xfrm')):
                off = xfrm.find(qn('a:off'))
                if off is not None:
                    x = int(off.get('x', '0')) - min_x
                    y = int(off.get('y', '0')) - min_y
                    off.set('x', str(x))
                    off.set('y', str(y))
            # 调整 a:cxn 连接点坐标
            for cxn in el.iter(qn('a:cxn')):
                x = int(cxn.get('x', '0')) - min_x
                y = int(cxn.get('y', '0')) - min_y
                cxn.set('x', str(x))
                cxn.set('y', str(y))
            # 移动到组中
            el.getparent().remove(el)
            grpSp.append(el)

        # 将组元素按正确顺序插入到 spTree 中
        # 放在原来第一个形状之后
        first_el = shapes[0]._element
        # first_el 已经被移除了，用它的兄弟元素作为参考
        spTree.append(grpSp)


def _group_all_content(slide):
    """
    将幻灯片中所有内容形状组合为一个组。
    这样连接线、横线、标签都在同一组内，移动任意元素时联动。
    用户在 PPT 中取消组合一次后即可单独编辑各元素。
    """
    from lxml import etree
    from pptx.oxml.ns import qn

    spTree = slide.shapes._spTree
    # 收集所有形状元素（排除 nvGrpSpPr/grpSpPr 等元数据）
    all_elements = []
    for child in list(spTree):
        tag = child.tag.split('}')[-1]
        if tag in ('sp', 'cxnSp', 'pic', 'grpSp'):
            all_elements.append(child)

    if len(all_elements) < 2:
        return

    # 计算包围盒
    min_x = min_y = None
    max_x = max_y = None
    for el in all_elements:
        for xfrm in el.iter(qn('a:xfrm')):
            off = xfrm.find(qn('a:off'))
            ext = xfrm.find(qn('a:ext'))
            if off is not None and ext is not None:
                x = int(off.get('x', '0'))
                y = int(off.get('y', '0'))
                cx = int(ext.get('cx', '0'))
                cy = int(ext.get('cy', '0'))
                if min_x is None or x < min_x:
                    min_x = x
                if min_y is None or y < min_y:
                    min_y = y
                if max_x is None or x + cx > max_x:
                    max_x = x + cx
                if max_y is None or y + cy > max_y:
                    max_y = y + cy
            # 对于连接器，检查 cxn 元素
            for cxn in xfrm.iter(qn('a:cxn')):
                x = int(cxn.get('x', '0'))
                y = int(cxn.get('y', '0'))
                parent_off = xfrm.find(qn('a:off'))
                if parent_off is not None:
                    x += int(parent_off.get('x', '0'))
                    y += int(parent_off.get('y', '0'))
                if min_x is None or x < min_x:
                    min_x = x
                if max_x is None or x > max_x:
                    max_x = x
                if min_y is None or y < min_y:
                    min_y = y
                if max_y is None or y > max_y:
                    max_y = y

    if min_x is None:
        return

    grp_w = max_x - min_x
    grp_h = max_y - min_y

    # 创建 <p:grpSp> 元素
    nsmap = {
        'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
        'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    }
    grpSp = etree.Element(qn('p:grpSp'), nsmap=nsmap)

    nv = etree.SubElement(grpSp, qn('p:nvGrpSpPr'))
    cnvPr = etree.SubElement(nv, qn('p:cNvPr'))
    cnvPr.set('id', '0')
    cnvPr.set('name', 'PES Group')
    etree.SubElement(nv, qn('p:cNvGrpSpPr'))
    etree.SubElement(nv, qn('p:nvPr'))

    gsp = etree.SubElement(grpSp, qn('p:grpSpPr'))
    gxfrm = etree.SubElement(gsp, qn('a:xfrm'))
    goff = etree.SubElement(gxfrm, qn('a:off'))
    goff.set('x', str(min_x))
    goff.set('y', str(min_y))
    gext = etree.SubElement(gxfrm, qn('a:ext'))
    gext.set('cx', str(grp_w))
    gext.set('cy', str(grp_h))
    chOff = etree.SubElement(gxfrm, qn('a:chOff'))
    chOff.set('x', '0')
    chOff.set('y', '0')
    chExt = etree.SubElement(gxfrm, qn('a:chExt'))
    chExt.set('cx', str(grp_w))
    chExt.set('cy', str(grp_h))

    # 移动所有形状到组，调整坐标
    for el in all_elements:
        for xfrm in el.iter(qn('a:xfrm')):
            off = xfrm.find(qn('a:off'))
            if off is not None:
                x = int(off.get('x', '0')) - min_x
                y = int(off.get('y', '0')) - min_y
                off.set('x', str(x))
                off.set('y', str(y))
        for cxn in el.iter(qn('a:cxn')):
            x = int(cxn.get('x', '0')) - min_x
            y = int(cxn.get('y', '0')) - min_y
            cxn.set('x', str(x))
            cxn.set('y', str(y))
        el.getparent().remove(el)
        grpSp.append(el)

    spTree.append(grpSp)


# 常量声明区 — 保持导入顺序



def hex_to_rgb(hex_str):
    """将 #RRGGBB 转换为 RGBColor"""
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def get_species_color(species):
    """获取物种颜色，未指定时按类型使用默认色"""
    if "color" in species and species["color"]:
        return hex_to_rgb(species["color"])
    stype = species.get("type", "well")
    if stype == "barrier":
        return hex_to_rgb(DEFAULT_BARRIER_COLOR)
    elif stype == "bimolecular":
        return hex_to_rgb(DEFAULT_BIMOL_COLOR)
    else:
        return hex_to_rgb(DEFAULT_WELL_COLOR)


# ============================================================
# 幻灯片尺寸 (16:9)
# ============================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN_L = Inches(1.2)
MARGIN_R = Inches(0.6)
MARGIN_T = Inches(1.0)   # 标题区留空
MARGIN_B = Inches(0.6)
PLOT_LEFT = MARGIN_L
PLOT_TOP = MARGIN_T
PLOT_WIDTH = SLIDE_WIDTH - MARGIN_L - MARGIN_R
PLOT_HEIGHT = Inches(4.8)  # 图表区高度
IMG_ROW_TOP = PLOT_TOP + PLOT_HEIGHT + Inches(0.4)  # 图片行顶部
IMG_ROW_HEIGHT = Inches(1.6)  # 图片行高度


# ============================================================
# 配置文件校验
# ============================================================
def validate_config(cfg):
    """校验 YAML 配置的必填字段"""
    errors = []
    if "title" not in cfg:
        errors.append("缺少 'title' 字段")
    if "species" not in cfg or not isinstance(cfg["species"], list) or len(cfg["species"]) < 2:
        errors.append("'species' 必须包含至少 2 个物种")

    for i, sp in enumerate(cfg.get("species", [])):
        if "name" not in sp:
            errors.append(f"species[{i}] 缺少 'name' 字段")
        if "energy" not in sp and "energy" not in sp:
            errors.append(f"species[{i}] 缺少 'energy' 字段")
        if sp.get("type") == "barrier" and "connects" not in sp:
            errors.append(f"species[{i}] (barrier) 缺少 'connects' 字段")

    if errors:
        print("配置校验失败:")
        for e in errors:
            print(f"  [ERR] {e}")
        sys.exit(1)
    return True


# ============================================================
# 布局计算
# ============================================================
def auto_positions(species):
    """为未指定 position 的物种自动分配横轴位置"""
    # 如果有 barrier 的 connects，优先按连接关系排序
    has_connects = any(s.get("type") == "barrier" for s in species)
    if not has_connects:
        names_seen = set()
        for i, sp in enumerate(species):
            if "position" not in sp:
                sp["position"] = i
            names_seen.add(sp["name"])
    else:
        # 构建连接图，进行拓扑排序
        # 简单策略：将所有 barriers 放在其连接的两个物种中间
        # 先确定哪些物种是 well/bimolecular
        wells = [s for s in species if s.get("type") != "barrier"]
        barriers = [s for s in species if s.get("type") == "barrier"]

        # 为 well 分配基础位置
        for i, w in enumerate(wells):
            if "position" not in w:
                w["position"] = i * 2  # 留出 barrier 的位置

        # 为 barrier 计算中间位置
        for b in barriers:
            if "position" not in b:
                conns = b.get("connects", [])
                positions = []
                for w in wells:
                    if w["name"] in conns and "position" in w:
                        positions.append(w["position"])
                if len(positions) >= 2:
                    b["position"] = (positions[0] + positions[1]) / 2.0
                elif len(positions) == 1:
                    b["position"] = positions[0] + 1.0
                else:
                    b["position"] = 0

    # 确保所有物种都有 position
    for sp in species:
        if "position" not in sp:
            sp["position"] = 0

    return species


def compute_axis_range(species, y_range=None):
    """计算坐标轴范围"""
    energies = [s.get("energy", 0) for s in species if "energy" in s]
    if not energies:
        return -5, 15
    min_e = min(energies)
    max_e = max(energies)
    padding = max(5, (max_e - min_e) * 0.15)
    if y_range:
        return y_range[0], y_range[1]
    return math.floor(min_e - padding - 1), math.ceil(max_e + padding + 1)


def data_to_screen(x_data, y_data, y_min, y_max, x_min, x_max):
    """将数据坐标转换为幻灯片上的像素坐标"""
    plot_left_px = PLOT_LEFT
    plot_top_px = PLOT_TOP
    plot_w_px = PLOT_WIDTH
    plot_h_px = PLOT_HEIGHT

    sx = plot_left_px + (x_data - x_min) / (x_max - x_min) * plot_w_px if x_max != x_min else plot_left_px + plot_w_px / 2
    sy = plot_top_px + plot_h_px - (y_data - y_min) / (y_max - y_min) * plot_h_px
    return int(sx), int(sy)


# ============================================================
# 绘图函数
# ============================================================
def _is_branched(species):
    """检测反应网络是否有分支（竞争反应）"""
    from collections import Counter
    barriers = [s for s in species if s.get("type") == "barrier"]
    if len(barriers) <= 1:
        return False
    # 统计每个物种被多少个 barrier 连接
    conn_count = Counter()
    for b in barriers:
        for c in b.get("connects", []):
            if c != b["name"]:
                conn_count[c] += 1
    # 如果有物种被 2+ 个 barrier 连接，说明有分支
    for name, count in conn_count.items():
        if count > 1:
            return True
    return False


def add_axes(slide, y_min, y_max, x_min, x_max, show_x_axis=True):
    """绘制坐标轴和刻度"""
    from pptx.util import Emu

    left_px = int(PLOT_LEFT)
    top_px = int(PLOT_TOP)
    right_px = int(PLOT_LEFT + PLOT_WIDTH)
    bottom_px = int(PLOT_TOP + PLOT_HEIGHT)

    # Y 轴位置：数据坐标 x=0 对应的屏幕位置
    axis_x, _ = data_to_screen(0, 0, y_min, y_max, x_min, x_max)
    axis_x = int(axis_x)

    if show_x_axis:
        # --- X 轴 (底部，从左侧边界到右侧边界) ---
        x_axis = slide.shapes.add_connector(
            1,  # MSO_CONNECTOR.STRAIGHT
            left_px, bottom_px, right_px, bottom_px
        )
        x_axis.line.color.rgb = hex_to_rgb(AXIS_COLOR)
        x_axis.line.width = Pt(1.5)
        _remove_shadow(x_axis)

        # X 轴箭头 (小三角形)
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.ISOSCELES_TRIANGLE,
            right_px - Inches(0.15), bottom_px - Inches(0.08),
            Inches(0.2), Inches(0.16)
        )
        arrow.rotation = 90.0
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = hex_to_rgb(AXIS_COLOR)
        arrow.line.fill.background()
        _remove_shadow(arrow)

        # X 轴标签
        x_label_box = slide.shapes.add_textbox(
            right_px - Inches(0.6), bottom_px + Inches(0.05),
            Inches(1.2), Inches(0.3)
        )
        x_label_tf = x_label_box.text_frame
        x_label_tf.text = "Reaction Coordinate"
        x_label_tf.paragraphs[0].font.size = Pt(11)
        x_label_tf.paragraphs[0].font.color.rgb = hex_to_rgb(AXIS_COLOR)
        x_label_tf.paragraphs[0].font.name = 'Times New Roman'
        x_label_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # --- Y 轴 (在 x=0 数据位置) ---
    y_axis = slide.shapes.add_connector(
        1,
        axis_x, top_px, axis_x, bottom_px
    )
    y_axis.line.color.rgb = hex_to_rgb(AXIS_COLOR)
    y_axis.line.width = Pt(1.5)
    _remove_shadow(y_axis)

    # Y 轴箭头
    arrow_y = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        axis_x - Inches(0.08), top_px - Inches(0.02),
        Inches(0.16), Inches(0.2)
    )
    arrow_y.fill.solid()
    arrow_y.fill.fore_color.rgb = hex_to_rgb(AXIS_COLOR)
    arrow_y.line.fill.background()
    _remove_shadow(arrow_y)

    # Y 轴标签 (紧贴 Y 轴顶部)
    y_label_box = slide.shapes.add_textbox(
        axis_x + Inches(0.05), top_px - Inches(0.05),
        Inches(1.0), Inches(0.3)
    )
    y_label_tf = y_label_box.text_frame
    y_label_tf.text = "E (kcal/mol)"
    y_label_tf.paragraphs[0].font.size = Pt(11)
    y_label_tf.paragraphs[0].font.color.rgb = hex_to_rgb(AXIS_COLOR)
    y_label_tf.paragraphs[0].font.name = 'Times New Roman'
    y_label_tf.paragraphs[0].alignment = PP_ALIGN.LEFT

    # --- Y 轴刻度 (对齐 Y 轴，标签紧贴纵轴) ---
    y_step = _nice_step(y_max - y_min)
    y_tick = math.ceil(y_min / y_step) * y_step
    while y_tick <= y_max:
        _, sy = data_to_screen(0, y_tick, y_min, y_max, x_min, x_max)
        # 刻度线 (从轴线向右延伸)
        tick_line = slide.shapes.add_connector(1, axis_x, sy, axis_x + Inches(0.08), sy)
        tick_line.line.color.rgb = hex_to_rgb(AXIS_COLOR)
        tick_line.line.width = Pt(1)
        _remove_shadow(tick_line)
        # 刻度标签 - 紧贴纵轴左侧
        tick_lbl = slide.shapes.add_textbox(axis_x - Inches(0.3), sy - Inches(0.1), Inches(0.28), Inches(0.2))
        tick_tf = tick_lbl.text_frame
        tick_tf.text = str(int(y_tick))
        tick_tf.paragraphs[0].font.size = Pt(9)
        tick_tf.paragraphs[0].font.color.rgb = hex_to_rgb(AXIS_COLOR)
        tick_tf.paragraphs[0].font.name = 'Times New Roman'
        tick_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        y_tick += y_step


def _nice_step(range_val):
    """计算刻度步长"""
    if range_val <= 0:
        return 5
    raw_step = range_val / 5.0
    magnitude = 10 ** math.floor(math.log10(raw_step))
    residual = raw_step / magnitude
    if residual <= 1.5:
        return magnitude
    elif residual <= 3.5:
        return 2 * magnitude
    elif residual <= 7.5:
        return 5 * magnitude
    else:
        return 10 * magnitude


def draw_energy_line(slide, species_list, y_min, y_max, x_min, x_max, marker_positions):
    """绘制连接的能量折线（从横线一端连到另一端，并连接到横线形状）"""
    barriers = [s for s in species_list if s.get("type") == "barrier"]
    wells = [s for s in species_list if s.get("type") != "barrier"]
    all_sp = {s["name"]: s for s in species_list}

    def _connect_endpoints(name_a, name_b, color):
        """从物种标记的端点到端点画连接线"""
        a_info = marker_positions.get(name_a)
        b_info = marker_positions.get(name_b)
        if a_info is None or b_info is None:
            return
        ax, ay, ahalf = a_info[:3]
        bx, by, bhalf = b_info[:3]
        # 左端的物种连右端 → 右端的物种连左端
        if ax < bx:
            x1, y1 = ax + ahalf, ay
            x2, y2 = bx - bhalf, by
        else:
            x1, y1 = ax - ahalf, ay
            x2, y2 = bx + bhalf, by
        _draw_line(slide, x1, y1, x2, y2, color)

    # 绘制 well → barrier → well 的连接
    drawn = set()
    for b in barriers:
        conns = b.get("connects", [])
        for c in conns:
            if c in all_sp and c != b["name"]:
                pair = tuple(sorted([b["name"], c]))
                if pair not in drawn:
                    drawn.add(pair)
                    color = get_species_color(all_sp[c])
                    _connect_endpoints(c, b["name"], color)

    # 如果是线性反应路径（无 barrier）
    if not barriers:
        sorted_wells = sorted(all_sp.values(), key=lambda s: s.get("position", 0))
        for i in range(len(sorted_wells) - 1):
            s1 = sorted_wells[i]
            s2 = sorted_wells[i + 1]
            pair = tuple(sorted([s1["name"], s2["name"]]))
            if pair not in drawn:
                drawn.add(pair)
                _connect_endpoints(s1["name"], s2["name"], get_species_color(s1))


def _draw_line(slide, x1, y1, x2, y2, color, width=None, dashed=True):
    """在幻灯片上画一条虚线（无阴影）"""
    if width is None:
        width = CONNECTOR_LINE_WIDTH
    connector = slide.shapes.add_connector(1, x1, y1, x2, y2)
    connector.line.color.rgb = color
    connector.line.width = width
    if dashed:
        connector.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    _remove_shadow(connector)


def draw_energy_point(slide, species, y_min, y_max, x_min, x_max):
    """在能量点处画横线标记和标签，返回 (中心x, 中心y, 半长)"""
    x_data = species.get("position", 0)
    y_data = species.get("energy", 0)
    sx, sy = data_to_screen(x_data, y_data, y_min, y_max, x_min, x_max)
    color = get_species_color(species)

    # 用横线代替点标记，长度默认 2cm
    half_len = int(SPECIES_MARKER_LENGTH / 2)
    marker = slide.shapes.add_connector(1, sx - half_len, sy, sx + half_len, sy)
    marker.line.color.rgb = color
    marker.line.width = SPECIES_MARKER_LINE_WIDTH
    _remove_shadow(marker)

    # 标签: 两行（名称 / 能量），居中对齐于横线下方
    label_w = Inches(1.2)
    label_h = Inches(0.5)
    label_top = sy  # 从横线高度开始向下
    label_box = slide.shapes.add_textbox(sx - label_w // 2, label_top, label_w, label_h)
    label_tf = label_box.text_frame
    label_tf.word_wrap = True

    # 第一行: 名称
    p = label_tf.paragraphs[0]
    p.text = species['name']
    p.font.size = Pt(11)
    p.font.color.rgb = color
    p.font.name = 'Times New Roman'
    p.alignment = PP_ALIGN.CENTER
    p.space_after = Pt(1)

    # 第二行: 能量值
    p2 = label_tf.add_paragraph()
    p2.text = f"{species.get('energy', 0):.1f}"
    p2.font.size = Pt(11)
    p2.font.color.rgb = color
    p2.font.name = 'Times New Roman'
    p2.alignment = PP_ALIGN.CENTER

    return sx, sy, half_len, marker, label_box  # 返回中心坐标、半长和所有形状引用


def embed_structure_images(slide, species_list, y_min, y_max, x_min, x_max, positions_cache):
    """在幻灯片底部嵌入分子结构图片，返回 {species_name: [pic_shape, name_label_shape]}"""
    with_images = [s for s in species_list if "image" in s and s["image"]]
    image_map = {}
    if not with_images:
        return image_map

    n = len(with_images)
    # 计算图片区域
    img_area_left = int(PLOT_LEFT)
    img_area_width = int(PLOT_WIDTH)
    img_area_top = int(IMG_ROW_TOP)

    # 每张图片可用宽度
    spacing = Inches(0.15)
    total_spacing = spacing * (n - 1) if n > 1 else 0
    img_w = int((img_area_width - total_spacing) / n)
    img_h = int(IMG_ROW_HEIGHT)

    for i, sp in enumerate(with_images):
        img_path = sp["image"]
        if not os.path.isfile(img_path):
            print(f"  [WARN] 图片文件不存在: {img_path}")
            continue

        img_left = img_area_left + i * (img_w + int(spacing))

        # 计算物种在图表中的水平位置对应的图片位置偏移
        x_data = sp.get("position", 0)
        sx, _ = data_to_screen(x_data, 0, y_min, y_max, x_min, x_max)
        # 如果 n 较大，让图片中心对齐物种在横轴的位置
        actual_left = int(sx - img_w / 2)
        # 限制边界
        actual_left = max(int(PLOT_LEFT), min(actual_left, int(PLOT_LEFT + PLOT_WIDTH - img_w)))

        try:
            pic = slide.shapes.add_picture(img_path, actual_left, img_area_top, img_w, img_h)
            # 给图片加边框
            pic.line.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
            pic.line.width = Pt(0.5)
        except Exception as e:
            print(f"  [WARN] 插入图片失败 {img_path}: {e}")
            continue

        # 图片下方标注物种名
        name_box = slide.shapes.add_textbox(
            actual_left, img_area_top + img_h + Inches(0.05),
            img_w, Inches(0.25)
        )
        name_tf = name_box.text_frame
        name_tf.word_wrap = True
        p = name_tf.paragraphs[0]
        p.text = sp["name"]
        p.font.size = Pt(9)
        p.font.color.rgb = get_species_color(sp)
        p.font.name = 'Times New Roman'
        p.alignment = PP_ALIGN.CENTER

        image_map[sp["name"]] = [pic, name_box]

    return image_map


# ============================================================
# PPT 生成核心函数
# ============================================================
def generate_pptx(cfg, output_path, mode="single"):
    """根据配置生成 PES 势能面 PPTX"""
    validate_config(cfg)

    title = cfg.get("title", "Potential Energy Surface")
    author = cfg.get("author", "")
    species = cfg.get("species", [])
    axis_cfg = cfg.get("axis", {})

    # 自动分配位置
    species = auto_positions(species)

    # 检测是否为分支反应（竞争反应），分支时强制隐藏 X 轴
    branched = _is_branched(species)
    show_x_axis = not branched and not axis_cfg.get("hide_x_axis", False)

    # 计算坐标轴范围
    y_range = axis_cfg.get("y_range", None)
    y_min, y_max = compute_axis_range(species, y_range)
    x_positions = [s.get("position", 0) for s in species]
    x_min = min(x_positions) - 0.3
    x_max = max(x_positions) + 0.3
    if x_max - x_min < 2:
        x_min -= 0.5
        x_max += 0.5

    prs = Presentation()
    prs.slide_width = int(SLIDE_WIDTH)
    prs.slide_height = int(SLIDE_HEIGHT)
    blank_layout = prs.slide_layouts[6]  # blank layout

    if mode == "single":
        # --- 标题幻灯片 ---
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title_slide(title_slide, title, author)

        # --- 内容幻灯片 ---
        content_slide = prs.slides.add_slide(blank_layout)
        add_slide_title(content_slide, title)
        add_axes(content_slide, y_min, y_max, x_min, x_max, show_x_axis=show_x_axis)

        # 先绘制物种横线标记，收集端点位置 + 形状引用
        marker_positions = {}
        species_group_map = {}
        for sp in species:
            sx, sy, half_len, marker, label_box = draw_energy_point(content_slide, sp, y_min, y_max, x_min, x_max)
            marker_positions[sp["name"]] = (sx, sy, half_len)
            species_group_map[sp["name"]] = [(marker, True), (label_box, False)]

        # 再绘制连接线（从横线一端到另一端，并粘合到横线形状）
        draw_energy_line(content_slide, species, y_min, y_max, x_min, x_max, marker_positions)

        # 嵌入分子结构图片，并入分组
        image_map = embed_structure_images(content_slide, species, y_min, y_max, x_min, x_max, marker_positions)
        for sp_name, img_shapes in image_map.items():
            if sp_name in species_group_map:
                species_group_map[sp_name].extend([(s, False) for s in img_shapes])

        # 每个物种组合为一个组
        _group_species_shapes(content_slide, species_group_map)
        # --- 标题幻灯片（stepwise 模式） ---
    elif mode == "stepwise":
        title_slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title_slide(title_slide, title, author)

        # 提取反应步骤：barrier 定义了步骤
        barriers = [s for s in species if s.get("type") == "barrier"]

        if not barriers:
            # 无 barrier，整张图
            content_slide = prs.slides.add_slide(blank_layout)
            add_slide_title(content_slide, title)
            add_axes(content_slide, y_min, y_max, x_min, x_max, show_x_axis=show_x_axis)
            marker_positions = {}
            species_group_map = {}
            for sp in species:
                sx, sy, half_len, marker, label_box = draw_energy_point(content_slide, sp, y_min, y_max, x_min, x_max)
                marker_positions[sp["name"]] = (sx, sy, half_len)
                species_group_map[sp["name"]] = [(marker, True), (label_box, False)]
            draw_energy_line(content_slide, species, y_min, y_max, x_min, x_max, marker_positions)
            image_map = embed_structure_images(content_slide, species, y_min, y_max, x_min, x_max, {})
            for sp_name, img_shapes in image_map.items():
                if sp_name in species_group_map:
                    species_group_map[sp_name].extend([(s, False) for s in img_shapes])
            _group_species_shapes(content_slide, species_group_map)
        else:
            # 分步展示
            name_map = {s["name"]: s for s in species}
            cumulative_names = set()
            for idx, b in enumerate(barriers):
                step_name = f"Step {idx + 1}: {b['name']}"
                conns = b.get("connects", [])
                step_species_names = set(conns) | {b["name"]}
                cumulative_names |= step_species_names

                step_slide = prs.slides.add_slide(blank_layout)
                add_slide_title(step_slide, f"{title} — {step_name}")

                # 只绘制到当前累积的物种
                step_species = [name_map[n] for n in cumulative_names if n in name_map]
                add_axes(step_slide, y_min, y_max, x_min, x_max, show_x_axis=show_x_axis)
                marker_positions = {}
                species_group_map = {}
                for sp in step_species:
                    sx, sy, half_len, marker, label_box = draw_energy_point(step_slide, sp, y_min, y_max, x_min, x_max)
                    marker_positions[sp["name"]] = (sx, sy, half_len)
                    species_group_map[sp["name"]] = [(marker, True), (label_box, False)]
                draw_energy_line(step_slide, step_species, y_min, y_max, x_min, x_max, marker_positions)
                image_map = embed_structure_images(step_slide, step_species, y_min, y_max, x_min, x_max, {})
                for sp_name, img_shapes in image_map.items():
                    if sp_name in species_group_map:
                        species_group_map[sp_name].extend([(s, False) for s in img_shapes])
                _group_species_shapes(step_slide, species_group_map)

            # 最终全景幻灯片
            final_slide = prs.slides.add_slide(blank_layout)
            add_slide_title(final_slide, f"{title} — 完整势能面")
            add_axes(final_slide, y_min, y_max, x_min, x_max, show_x_axis=show_x_axis)
            marker_positions = {}
            species_group_map = {}
            for sp in species:
                sx, sy, half_len, marker, label_box = draw_energy_point(final_slide, sp, y_min, y_max, x_min, x_max)
                marker_positions[sp["name"]] = (sx, sy, half_len)
                species_group_map[sp["name"]] = [(marker, True), (label_box, False)]
            draw_energy_line(final_slide, species, y_min, y_max, x_min, x_max, marker_positions)
            image_map = embed_structure_images(final_slide, species, y_min, y_max, x_min, x_max, {})
            for sp_name, img_shapes in image_map.items():
                if sp_name in species_group_map:
                    species_group_map[sp_name].extend([(s, False) for s in img_shapes])
            _group_species_shapes(final_slide, species_group_map)

    # 保存
    prs.save(output_path)
    print(f"[OK] 已生成 PES 势能面 PPTX: {output_path}")
    print(f"   幻灯片数: {len(prs.slides)}")


def add_title_slide(slide, title, author=""):
    """添加标题幻灯片"""
    # 标题
    title_box = slide.shapes.add_textbox(
        Inches(1.5), Inches(2.0), Inches(10), Inches(1.5)
    )
    tf = title_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.color.rgb = hex_to_rgb("#000000")
    p.font.name = 'Times New Roman'
    p.alignment = PP_ALIGN.CENTER

    if author:
        author_box = slide.shapes.add_textbox(
            Inches(1.5), Inches(3.8), Inches(10), Inches(0.5)
        )
        tf2 = author_box.text_frame
        tf2.paragraphs[0].text = author
        tf2.paragraphs[0].font.size = Pt(16)
        tf2.paragraphs[0].font.color.rgb = RGBColor(0x66, 0x66, 0x66)
        tf2.paragraphs[0].font.name = 'Times New Roman'
        tf2.paragraphs[0].alignment = PP_ALIGN.CENTER


def add_slide_title(slide, title):
    """在内容幻灯片的顶部添加标题"""
    title_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15), Inches(12), Inches(0.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(20)
    p.font.color.rgb = hex_to_rgb("#000000")
    p.font.name = 'Times New Roman'
    p.alignment = PP_ALIGN.LEFT


# ============================================================
# CLI 入口
# ============================================================
def main():
    parser = argparse.ArgumentParser(description="PES Sketch — 生成反应势能面 PPTX 草图")
    parser.add_argument("config", help="YAML 配置文件路径")
    parser.add_argument("-o", "--output", default="pes_sketch.pptx", help="输出 PPTX 文件路径")
    parser.add_argument("--mode", choices=["single", "stepwise"], default="single",
                        help="幻灯片模式: single=单页全景, stepwise=分步展示")
    args = parser.parse_args()

    if not os.path.isfile(args.config):
        print(f"[ERR] 配置文件不存在: {args.config}")
        sys.exit(1)

    with open(args.config, "r", encoding="utf-8") as f:
        cfg = yaml.safe_load(f)

    generate_pptx(cfg, args.output, mode=args.mode)


if __name__ == "__main__":
    main()
